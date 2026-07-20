"""Вставить фото в плейсхолдер «фото» презентации для ИТМО.

    python scripts/put_photo_itmo.py                      # показать, куда что можно вставить
    python scripts/put_photo_itmo.py 1 look-editorial.png # слайд 1 <- картинка

Зачем скрипт, а не «перетащить мышкой»: плейсхолдеры в фирменном шаблоне разного формата, и картинка,
вставленная как есть, растягивается. Здесь она вписывается по принципу «заполнить»: масштабируется по
короткой стороне и обрезается по длинной, пропорции лица и фигуры не плывут.

Правки, сделанные руками в PowerPoint, скрипт не трогает — он работает по месту плейсхолдера.
Пересборка колоды (build_submission_pptx_itmo.py) вставленные фото ЗАТРЁТ: сначала текст, потом фото.
"""
from pathlib import Path
import sys

from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent.parent
DECK = ROOT / "submission" / "pptx" / "02-презентация-ИТМО.pptx"


def placeholders(prs):
    """Фигуры-плейсхолдеры «фото»: (номер слайда, фигура)."""
    found = []
    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip().replace("​", "") == "фото":
                found.append((idx, shape))
    return found


def fill_crop(img_path: Path, box_w: int, box_h: int):
    """Коэффициенты обрезки, чтобы картинка заполнила рамку без искажения пропорций."""
    from PIL import Image
    with Image.open(img_path) as im:
        iw, ih = im.size
    box_ratio = box_w / box_h
    img_ratio = iw / ih
    if img_ratio > box_ratio:          # картинка шире рамки — режем бока
        keep = box_ratio / img_ratio
        return (1 - keep) / 2, 0.0
    keep = img_ratio / box_ratio       # картинка выше рамки — режем верх и низ
    return 0.0, (1 - keep) / 2


def main() -> int:
    if not DECK.exists():
        print(f"Нет колоды: {DECK.relative_to(ROOT)} — сначала собери build_submission_pptx_itmo.py")
        return 1

    prs = Presentation(str(DECK))
    spots = placeholders(prs)

    if len(sys.argv) < 3:
        print(f"Свободных плейсхолдеров «фото»: {len(spots)}")
        for num, shape in spots:
            # Знак умножения не пишем: консоль Windows в cp1251 роняет его UnicodeEncodeError.
            print(f"  слайд {num}: рамка {shape.width / 914400:.1f}x{shape.height / 914400:.1f} дюйма"
                  f" ({'вертикальная' if shape.height > shape.width else 'горизонтальная'})")
        print("\nВставить:  python scripts/put_photo_itmo.py <номер слайда> <файл>")
        return 0

    slide_no = int(sys.argv[1])
    img = Path(sys.argv[2])
    if not img.is_absolute():
        img = ROOT / img
    if not img.exists():
        print(f"Нет файла: {img}")
        return 1

    target = [(n, s) for n, s in spots if n == slide_no]
    if not target:
        print(f"На слайде {slide_no} нет свободного плейсхолдера «фото». Занятые не перезаписываю.")
        return 1

    _, shape = target[0]
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    slide = prs.slides[slide_no - 1]

    pic = slide.shapes.add_picture(str(img), left, top, width, height)
    crop_x, crop_y = fill_crop(img, width, height)
    pic.crop_left = pic.crop_right = crop_x
    pic.crop_top = pic.crop_bottom = crop_y

    # Картинку — на место плейсхолдера в порядке отрисовки, сам плейсхолдер убираем.
    shape._element.addprevious(pic._element)
    shape._element.getparent().remove(shape._element)

    prs.save(str(DECK))
    print(f"OK: слайд {slide_no} <- {img.name} (обрезка {crop_x:.0%}/{crop_y:.0%})")
    print("Проверь в PowerPoint: кадрирование выбрано по центру, при необходимости подвинь.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
