# -*- coding: utf-8 -*-
"""Сборка БРЕНДОВОЙ презентации подачи ИТМО (.pptx) в стиле Petrogradka Editorial.

Зачем отдельно от build_submission_pptx_itmo.py: тот заполняет казённый шаблон ИТМО (синий,
корпоративный). Этот собирает презентацию в фирменном стиле проекта — кремовый фон, винный акцент,
сериф-заголовки, много воздуха. Положение о конкурсе (§2.2) просит «дополнительные презентационные
материалы» без обязательного шаблона, поэтому брендовая версия допустима.

Контент выверен по ПЯТИ критериям ИТМО (§3.4.1 Положения):
  1. Разработка и инженерия — слайды «Пайплайн» и числа инженерии
  2. Data Science          — слайд eval против экспертной разметки
  3. Применение ИИ         — слайд «Как построен проект» (агентная разработка)
  4. Продуктовое мышление  — проблема, ЦА, конкуренты, MVP, обратная связь
  5. Мотивация             — финальный слайд

Шрифты: Georgia (сериф с кириллицей, есть на любом Windows) вместо брендового Cormorant Garamond —
Cormorant в системе не установлен, PowerPoint молча подменил бы его на Calibri и сломал вёрстку на
чужой машине/проекторе. Палитра — Petrogradka Editorial из architecture/visual-direction.md.

Реплики докладчика кладём в заметки к слайдам (view: Заметки) — их видно в режиме докладчика.

ЗАПУСК:
    python scripts/build_submission_pptx.py
    python scripts/build_submission_pptx.py --out submission/pptx/prezentaciya.pptx
"""
from __future__ import annotations
import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT_DEFAULT = ROOT / "submission" / "pptx" / "02-презентация.pptx"

# Petrogradka Editorial
CREAM = RGBColor(0xF5, 0xEF, 0xE3)
INK = RGBColor(0x1F, 0x1D, 0x1B)
WINE = RGBColor(0x5D, 0x22, 0x30)
MUTED = RGBColor(0x6B, 0x64, 0x5C)
LINE = RGBColor(0xE3, 0xDC, 0xCF)

SERIF = "Georgia"       # заголовки (брендовый Cormorant не установлен — см. докстринг)
SANS = "Segoe UI"       # текст


def _bg(slide, prs):
    """Кремовый фон — на всю площадь слайда."""
    from pptx.enum.shapes import MSO_SHAPE
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid()
    r.fill.fore_color.rgb = CREAM
    r.line.fill.background()
    r.shadow.inherit = False
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)   # фон — под остальным


def _text(slide, left, top, width, height, text, size, font=SANS, color=INK,
          bold=False, align=PP_ALIGN.LEFT, italic=False, spacing=1.15):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        run = p.add_run()
        run.text = ln
        run.font.size = Pt(size)
        run.font.name = font
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def _rule(slide, left, top, width):
    """Тонкая винная линия — тот же приём, что в вёрстке Карты."""
    from pptx.enum.shapes import MSO_SHAPE
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(2))
    r.fill.solid()
    r.fill.fore_color.rgb = WINE
    r.line.fill.background()
    r.shadow.inherit = False


def _eyebrow(slide, text):
    _text(slide, 0.9, 0.5, 8, 0.3, text.upper(), 11, SANS, WINE, bold=True)
    _rule(slide, 0.9, 0.85, 11.5)


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# Кадры «до/после» с лендинга: те же, что уже опубликованы на сайте (фаундер, согласие есть).
# Эмоциональный якорь презентации — по заметкам к контенту повторяются на титуле и в результатах.
PHOTO_BEFORE = ROOT / "web" / "photos" / "hero" / "01-do.jpg"
PHOTO_AFTER = ROOT / "web" / "photos" / "hero" / "02-posle.jpg"


def _pair_photos(slide, left, top, height, caption_size=11, slide_w=13.333):
    """Два кадра «до/после» рядом + подписи. Кадры вертикальные 3:4 — ширину считаем от высоты.

    Высоту ужимаем, если пара не влезает в остаток слайда: два вертикальных кадра съедают ширины
    больше, чем кажется (h=4.9" → 2×3.68" + зазор = 7.5"), и молча уезжают за край.
    """
    gap = 0.18
    avail = slide_w - left - 0.35          # правое поле
    width = height * 0.75
    if 2 * width + gap > avail:
        width = (avail - gap) / 2
        height = width / 0.75
    for i, (path, cap) in enumerate([(PHOTO_BEFORE, "до"), (PHOTO_AFTER, "после")]):
        if not path.exists():
            continue
        x = left + i * (width + gap)
        slide.shapes.add_picture(str(path), Inches(x), Inches(top),
                                 width=Inches(width), height=Inches(height))
        _text(slide, x, top + height + 0.04, width, 0.3, cap.upper(), caption_size, SANS,
              WINE if cap == "после" else MUTED, bold=True, align=PP_ALIGN.CENTER)


def _bullets(slide, items, top=2.3, size=17, gap=0.62, width=11.0):
    for i, it in enumerate(items):
        _text(slide, 1.15, top + i * gap, width, 0.55, "— " + it, size, SANS, INK)


def build(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]

    # ── 1. Титул ───────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, prs)
    _text(s, 0.9, 2.0, 7.4, 1.2, "Sense Style AI", 54, SERIF, INK)
    _rule(s, 0.9, 3.25, 4.2)
    _text(s, 0.9, 3.5, 7.4, 0.9, "Персональный стилист\nна основе психологии моды", 22, SERIF, WINE, italic=True)
    _text(s, 0.9, 4.75, 7.4, 1.2,
          "Измеряем разрыв между тем, какой женщину видят сейчас,\nи тем, какой она хочет быть — и закрываем его гардеробом.",
          15, SANS, MUTED)
    _text(s, 0.9, 6.4, 7.4, 0.5, "Ксения Колупаева  ·  ИИ-стартап  ·  Junior ML Contest ИТМО, 2026", 13, SANS, MUTED)
    _pair_photos(s, left=7.5, top=1.6, height=4.4)
    _notes(s, "0:00–0:20\n\n«Sense Style AI измеряет разрыв между тем, какой женщину видят сейчас, "
              "и тем, какой она хочет быть, — и собирает гардероб, который этот разрыв закрывает.»\n\n"
              "ВИЗУАЛ: кадры «до/после» — те же, что на лендинге. Эмоциональный якорь, повторяется "
              "на слайде 7.")

    # ── 2. Проблема + ЦА ───────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, prs)
    _eyebrow(s, "01 · Проблема и аудитория")
    _text(s, 0.9, 1.15, 11.5, 0.9, "«Полный шкаф, а надеть нечего»", 36, SERIF, INK)
    _bullets(s, [
        "ЦА: женщины 30–50 в точке перехода — новая должность, материнство, переезд",
        "Внутренне изменилась, а образ остался прежним — не выглядит как та, кем стала",
        "Симптом «нечего надеть» — не про вещи, а про разрыв идентичности",
    ])
    _text(s, 1.15, 4.7, 11, 0.8, "AI-стилисты советуют вещи по трендам. Мы измеряем разрыв — и закрываем именно его.",
          20, SERIF, WINE, italic=True)
    _text(s, 1.15, 5.9, 11, 0.5, "Существующие AI-стилисты решают задачу рекомендации. Наша задача другая — и измеримая.",
          14, SANS, MUTED)
    _notes(s, "0:20–1:00\n\n«Женщина в точке перехода — новая должность, материнство, переезд — внутренне "
              "изменилась, а образ остался прежним. Она не выглядит как та, кем стала. Существующие "
              "AI-стилисты советуют вещи по трендам. Это не решает проблему: проблема не в вещах, "
              "а в разрыве идентичности.»")

    # ── 3. Решение, наука, конкуренты ──────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, prs)
    _eyebrow(s, "02 · Решение · наука · рынок")
    _text(s, 0.9, 1.1, 11.5, 0.9, "Identity Gap — разрыв в процентах", 34, SERIF, INK)
    _text(s, 0.9, 2.05, 5.6, 1.4, "46–78%", 72, SERIF, WINE)
    _text(s, 0.9, 3.45, 5.6, 0.9, "разброс стартового разрыва\nна живых клиентках — шкала различает состояния",
          13, SANS, MUTED)
    # Научная опора
    _text(s, 0.9, 4.7, 5.6, 0.4, "Научная опора", 15, SERIF, INK, bold=True)
    _text(s, 0.9, 5.15, 5.6, 0.4, "Self-Discrepancy Theory · Higgins, 1987", 13, SANS, WINE, bold=True)
    _text(s, 0.9, 5.5, 5.6, 0.4, "Enclothed Cognition · Adam & Galinsky, 2012", 13, SANS, WINE, bold=True)
    _text(s, 0.9, 5.95, 5.6, 0.8, "Поверх — авторская методология: 4 стиля, 25 подстилей, 7-шаговый алгоритм. Она и есть наша разметка данных.",
          12, SANS, MUTED)
    # Конкуренты — правая колонка (продуктовое мышление)
    _text(s, 6.9, 2.05, 5.6, 0.4, "Чем мы отличаемся от рынка", 15, SERIF, INK, bold=True)
    comp = [
        ("Персональный шопинг (Stitch Fix)", "оптимизирует покупку, не самовосприятие"),
        ("Гардеробы (Whering, Acloset)", "собирают из того, что есть, без вопроса «кем хочешь быть»"),
        ("Цветотип по фото (Style DNA)", "колористика в отрыве от запроса клиентки"),
    ]
    top = 2.6
    for name, sub in comp:
        _text(s, 6.9, top, 5.6, 0.4, name, 13, SANS, WINE, bold=True)
        _text(s, 6.9, top + 0.36, 5.6, 0.5, sub, 12, SANS, MUTED)
        top += 1.02
    _text(s, 6.9, 5.75, 5.6, 1.0, "Общее у всех — результат не измеряется. Наша гипотеза: разрыв, выраженный числом, делает ценность проверяемой.",
          13, SERIF, WINE, italic=True)
    _notes(s, "1:00–1:45\n\n«Мы операционализировали разрыв идентичности в процент — Identity Gap. Опора "
              "научная: Self-Discrepancy Theory и Enclothed Cognition. Поверх — авторская методология, "
              "она и есть наша разметка данных. От рынка отличаемся тем, что у конкурентов результат не "
              "измеряется — а мы выражаем разрыв числом и показываем динамику.»")

    # ── 4. Пайплайн (Разработка и инженерия) ───────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, prs)
    _eyebrow(s, "03 · Как это работает")
    _text(s, 0.9, 1.15, 11.5, 0.9, "Пайплайн: от анкеты до образа на клиентке", 34, SERIF, INK)
    steps = [
        ("Квиз + фото", "мультимодальный вход"),
        ("Vision", "колорит, геометрия фигуры"),
        ("RAG", "правила метода, объяснимость"),
        ("Диагностика", "Формула стиля + Identity Gap"),
        ("Генерация", "образы на её лице и фигуре"),
    ]
    x = 0.9
    for i, (title, sub) in enumerate(steps):
        _text(s, x, 2.6, 2.1, 0.4, title, 15, SANS, WINE, bold=True)
        _text(s, x, 3.05, 2.1, 0.8, sub, 12, SANS, MUTED)
        if i < len(steps) - 1:
            _text(s, x + 1.95, 2.55, 0.4, 0.4, "→", 20, SANS, WINE)
        x += 2.42
    _text(s, 0.9, 4.6, 11.5, 1.4,
          "Ключевое инженерное решение: Identity Gap считается ОДИН раз на сервере — и одинаково виден\n"
          "на всех экранах. Не бывает «на квизе одно число, на Карте другое».",
          17, SANS, INK)
    _text(s, 0.9, 5.9, 11.5, 0.6, "Полная Карта стиля собирается за ~66 секунд (живой прогон на проде).",
          15, SERIF, WINE, italic=True)
    _notes(s, "1:45–2:30\n\n«Вход мультимодальный: анкета и фото. Vision читает колорит и геометрию фигуры. "
              "RAG подмешивает правила метода и объясняет диагноз. Диагностика выдаёт Формулу стиля и "
              "Identity Gap. Генератор примеряет образы на её собственное лицо и фигуру. Ключевое: Gap "
              "считается один раз на сервере и одинаково виден на всех экранах.»\n\n"
              "Стек: Vision + генерация через OpenRouter, провайдер сменяем за один модуль. Интеллект — "
              "в версионируемой промпт-библиотеке со строгим JSON.")

    # ── 5. Применение ИИ = агентная разработка ─────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, prs)
    _eyebrow(s, "04 · Применение ИИ — как построен проект")
    _text(s, 0.9, 1.15, 11.5, 0.9, "Проект построен в агентной разработке", 34, SERIF, INK)
    _text(s, 0.9, 2.15, 6.0, 1.4, "254 / 263", 60, SERIF, WINE)
    _text(s, 0.9, 3.5, 6.0, 0.7, "коммита — в соавторстве с кодовым агентом.\nМетодология и решения — мои, агент вёл реализацию.",
          14, SANS, MUTED)
    cards = [
        ("Регламент для агента", "AGENTS.md и CLAUDE.md: источники правды, границы правок, план в plans/ на каждую функцию"),
        ("Тесты как контур", "328 автотестов держат агента: регресс закрывается тестом, а не правкой на глаз"),
        ("Агент в данных и eval", "индексация методологии в RAG, разбор кейсов в ground truth, ablation с продуктовой гипотезой"),
    ]
    top = 2.25
    for name, body in cards:
        _text(s, 6.9, top, 5.6, 0.4, name, 15, SANS, WINE, bold=True)
        _text(s, 6.9, top + 0.38, 5.6, 0.9, body, 12, SANS, INK)
        top += 1.32
    _text(s, 0.9, 5.0, 6.0, 1.4, "Критерий «Применение ИИ» — про AI-агентов в работе над проектом. Продуктовый стек (Vision, RAG, генерация) — на слайдах 3 и 6.",
          12, SERIF, WINE, italic=True)
    _notes(s, "2:30–3:10\n\n«Проект целиком построен в агентной разработке: 254 из 263 коммитов — в "
              "соавторстве с кодовым агентом. Это процесс с правилами: AGENTS.md и CLAUDE.md как "
              "регламент, 328 тестов как страховка, агент в работе с данными и в исследовании. "
              "Постановка задачи, методология и решения — мои; агент — сильный исполнитель.»\n\n"
              "Если спросят «что делали сами» — методология, архитектура, критерии качества, выбор "
              "моделей. Агент не знает, что натуральной зиме не идёт беж — это из методологии.")

    # ── 6. Data Science ────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, prs)
    _eyebrow(s, "05 · Data Science: оценка")
    _text(s, 0.9, 1.15, 11.5, 0.9, "Меряем диагностику против экспертной разметки", 32, SERIF, INK)

    rows, cols = 3, 4
    tbl = s.shapes.add_table(rows, cols, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.5)).table
    headers = ["Вход диагностики", "dominant_field_accuracy", "formula_hit_rate", "gap_sanity"]
    data = [["3 пика (как в проде)", "0.667", "1.0", "1.0"],
            ["полный набор черт (ablation)", "1.000", "1.0", "1.0"]]
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.size, r.font.name, r.font.bold = Pt(12), SANS, True
                r.font.color.rgb = CREAM
        c.fill.solid()
        c.fill.fore_color.rgb = WINE
    for i, row in enumerate(data, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = val
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size, r.font.name = Pt(13), SANS
                    r.font.color.rgb = INK
                    r.font.bold = (i == 2 and j == 1)
            c.fill.solid()
            c.fill.fore_color.rgb = CREAM

    _text(s, 0.9, 4.2, 11.5, 1.0,
          "Инсайт из данных: продукт обрезает желаемые черты до трёх и теряет сигнал —\n"
          "на полном наборе попадание в поле эксперта растёт с 2/3 до 3/3.",
          18, SANS, INK)
    _text(s, 0.9, 5.5, 11.5, 0.8,
          "Честно: выборка мала (n = 3). Vision изолирован, чтобы мерить именно диагностику.\n"
          "Метрику расширяем по мере накопления кейсов.",
          14, SERIF, WINE, italic=True)
    _notes(s, "3:10–3:55\n\n«Мы оцениваем диагностику против экспертной разметки. Vision изолируем, чтобы "
              "мерить именно диагноз. И сразу получили инсайт из данных: продукт обрезает желаемые черты "
              "до трёх и теряет сигнал — на полном наборе точность попадания в поле эксперта растёт "
              "с 2/3 до 3/3. Оговорюсь честно: выборка пока мала, n=3, метрику расширяем.»\n\n"
              "НЕ проговаривать все числа — показать таблицу, назвать один вывод (2/3 → 3/3).")

    # ── 7. Результаты + обратная связь ─────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, prs)
    _eyebrow(s, "06 · Результаты и обратная связь")
    _text(s, 0.9, 1.15, 8.6, 0.9, "MVP работает end-to-end —\nна реальных клиентках", 30, SERIF, INK)
    # Числа — только по живым клиенткам: смоук-тесты и самотесты автора из воронки исключены
    # (иначе меряем свою же работу; жюри разбирает такую конверсию первым вопросом).
    kpis = [("83%", "конверсия\nквиз → Карта"), ("46–78%", "разброс\nIdentity Gap"),
            ("66 с", "полная Карта\nс образами"), ("328", "тестов,\nCI зелёный")]
    x = 0.9
    for val, sub in kpis:
        _text(s, x, 2.55, 2.0, 0.8, val, 30, SERIF, WINE)
        _text(s, x, 3.4, 2.0, 0.8, sub, 12, SANS, MUTED)
        x += 2.15
    _bullets(s, [
        "Метод валидирован двумя профильными экспертами, а не только автором",
        "Прогон клиентки поймал критический баг → фикс → регрессионный тест",
        "Две клиентки независимо указали на монотонность палитры → правила качества",
        "Публичный репозиторий, Docker, CI со сборкой образа и smoke-тестом",
    ], top=4.35, size=13, gap=0.42, width=7.7)
    _text(s, 0.9, 6.2, 7.9, 0.8,
          "Ценность меряем не «понравилось», а процентом закрытия Identity Gap.\nЛонгитюд «до/после» набираем: эффект от ношения требует недель.",
          12, SERIF, WINE, italic=True)
    _pair_photos(s, left=8.9, top=1.9, height=4.0, caption_size=10)
    _notes(s, "3:55–4:40\n\n«MVP работает end-to-end: полная Карта за минуту. Конверсия квиз→Карта 83%. "
              "Метод валидирован двумя профильными экспертами. Обратная связь меняла продукт: прогон "
              "клиентки поймал критический баг, две клиентки независимо указали на монотонность палитры. "
              "Инженерия закрыта — репозиторий, Docker, CI, 328 тестов.»\n\n"
              "ЧИСЛА ЧЕСТНЫЕ: считаем живых клиенток, смоук-тесты и самотесты автора из воронки исключены.\n\n"
              "ВАЖНО про «до/после»: НЕ заявлять, что трансформация измерена на дистанции. Лонгитюда нет — "
              "честная формулировка: «механика замера вшита, инфраструктура готова, лонгитюд требует "
              "недель». Работает как оговорка про n=3.")

    # ── 8. Мотивация ───────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _bg(s, prs)
    _eyebrow(s, "07 · Мотивация")
    _text(s, 0.9, 1.5, 11.5, 1.6,
          "Я построила эту систему как доменный эксперт,\nкоторый сам собрал AI-слой.",
          32, SERIF, INK)
    _bullets(s, [
        "15+ лет в авиации: инженерно-навигационные расчёты — привычка к точности",
        "Имидж-консалтинг: авторская методология, валидированная профильными экспертами",
        "AITH: выпускную веду как MVP стартапа, роль Product Engineer — ровно моя позиция",
        "Цель: довести Data Science и инженерию до продакшн-качества",
    ], top=3.6, size=17, gap=0.62)
    _text(s, 0.9, 6.3, 11.5, 0.6, "Спасибо.", 24, SERIF, WINE, italic=True)
    _notes(s, "4:40–5:00\n\n«Я построила эту систему как доменный эксперт, который сам собрал AI-слой. "
              "В AITH выпускную можно вести как MVP стартапа, а роль Product Engineer описывает ровно "
              "мою позицию. Хочу углубить Data Science и инженерию, чтобы довести оценку и генерацию "
              "до продакшн-качества. Спасибо.»\n\nДержать 5:00 — на вопросы отдельные 7 минут "
              "(см. 03-сценарий-защиты.md и 09-вопросы-по-цифрам.md).")


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=str(OUT_DEFAULT))
    a = ap.parse_args()

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    build(prs)

    out = Path(a.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    # маркеры ASCII: консоль Windows в cp1251 роняет «✓»/«→» уже ПОСЛЕ записи файла
    print(f"OK: {len(prs.slides.__iter__.__self__._sldIdLst)} слайдов -> {out.relative_to(ROOT)}")
    print("Открой в PowerPoint: правь текст, при желании замени кадры «до/после» на слайдах 1 и 7,")
    print("затем Файл -> Сохранить как -> PDF.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
